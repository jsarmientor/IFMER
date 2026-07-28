import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images(pptx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_filename = f"slide_{slide_num}_img_{j}.{image_ext}"
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                print(f"Saved {image_path}")

if __name__ == "__main__":
    if len(sys.argv) > 2:
        extract_images(sys.argv[1], sys.argv[2])
    else:
        print("Usage: python extract_images.py <pptx_path> <output_dir>")
