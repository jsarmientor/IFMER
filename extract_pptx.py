import sys
import collections
import collections.abc

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is required.")
    sys.exit(1)

def extract_text(file_path, output_path):
    try:
        prs = Presentation(file_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            for i, slide in enumerate(prs.slides):
                f.write(f"\n--- Slide {i+1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        f.write(shape.text.strip() + "\n")
    except Exception as e:
        print(f"Error reading presentation: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 2:
        extract_text(sys.argv[1], sys.argv[2])
    else:
        print("Usage: python extract_pptx.py <input> <output>")
